#!/usr/bin/env python3
"""
HTML2PPT 快速开始脚本
自动检测环境并选择最佳安装方式
"""

import subprocess
import sys
import os
import platform

def check_command_exists(command):
    """检查命令是否存在"""
    try:
        subprocess.run([command, "--version"], capture_output=True, check=True)
        return True
    except (subprocess.CalledProcessError, FileNotFoundError):
        return False

def detect_environment():
    """检测运行环境"""
    env_info = {
        'os': platform.system(),
        'python_version': f"{sys.version_info.major}.{sys.version_info.minor}",
        'has_conda': check_command_exists('conda'),
        'has_pip': check_command_exists('pip'),
        'in_venv': hasattr(sys, 'real_prefix') or (hasattr(sys, 'base_prefix') and sys.base_prefix != sys.prefix)
    }
    return env_info

def print_environment_info(env_info):
    """打印环境信息"""
    print("🔍 环境检测结果:")
    print(f"  操作系统: {env_info['os']}")
    print(f"  Python版本: {env_info['python_version']}")
    print(f"  Conda可用: {'✅' if env_info['has_conda'] else '❌'}")
    print(f"  pip可用: {'✅' if env_info['has_pip'] else '❌'}")
    print(f"  虚拟环境: {'✅' if env_info['in_venv'] else '❌'}")
    print()

def recommend_installation(env_info):
    """推荐安装方式"""
    if not env_info['has_pip']:
        return "error", "pip不可用，请先安装pip"
    
    if float(env_info['python_version']) < 3.8:
        return "error", f"Python版本过低({env_info['python_version']})，需要3.8或更高版本"
    
    if env_info['has_conda'] and not env_info['in_venv']:
        return "conda", "推荐使用conda创建环境并安装"
    elif not env_info['in_venv']:
        return "venv", "推荐创建虚拟环境后安装"
    else:
        return "direct", "可以直接安装"

def run_installation(method, env_info):
    """执行安装"""
    print(f"🚀 开始安装 (方法: {method})")
    
    if method == "conda":
        print("使用conda安装...")
        commands = [
            "conda install -c conda-forge python-pptx beautifulsoup4 lxml pillow requests pyyaml -y",
            f"{sys.executable} -m pip install cssutils html2text markdown"
        ]
    elif method == "venv":
        print("创建虚拟环境并安装...")
        venv_name = "html2ppt_env"
        if env_info['os'] == 'Windows':
            activate_cmd = f"{venv_name}\\Scripts\\activate"
            pip_cmd = f"{venv_name}\\Scripts\\pip"
        else:
            activate_cmd = f"source {venv_name}/bin/activate"
            pip_cmd = f"{venv_name}/bin/pip"
        
        commands = [
            f"{sys.executable} -m venv {venv_name}",
            f"{pip_cmd} install --upgrade pip",
            f"{pip_cmd} install -r requirements-minimal.txt"
        ]
        print(f"虚拟环境将创建在: {venv_name}")
        print(f"激活命令: {activate_cmd}")
    else:  # direct
        print("直接安装...")
        commands = [
            f"{sys.executable} install.py"
        ]
    
    success = True
    for cmd in commands:
        print(f"执行: {cmd}")
        try:
            result = subprocess.run(cmd, shell=True, check=True)
            print("✅ 成功")
        except subprocess.CalledProcessError as e:
            print(f"❌ 失败: {e}")
            success = False
            break
    
    return success

def test_installation():
    """测试安装结果"""
    print("\n🧪 测试安装...")
    
    test_script = """
try:
    import pptx
    from bs4 import BeautifulSoup
    import lxml
    from PIL import Image
    import requests
    import yaml
    print("✅ 核心依赖导入成功")
    
    # 简单功能测试
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    print("✅ PPT创建功能正常")
    
    soup = BeautifulSoup("<html><body><h1>Test</h1></body></html>", 'lxml')
    print("✅ HTML解析功能正常")
    
    print("🎉 安装测试通过！")
    
except ImportError as e:
    print(f"❌ 导入失败: {e}")
    exit(1)
except Exception as e:
    print(f"❌ 功能测试失败: {e}")
    exit(1)
"""
    
    try:
        result = subprocess.run([sys.executable, "-c", test_script], 
                              capture_output=True, text=True, check=True)
        print(result.stdout)
        return True
    except subprocess.CalledProcessError as e:
        print("❌ 测试失败")
        if e.stdout:
            print("输出:", e.stdout)
        if e.stderr:
            print("错误:", e.stderr)
        return False

def show_next_steps():
    """显示后续步骤"""
    print("\n📖 后续步骤:")
    print("1. 运行示例:")
    print("   python run_example.py")
    print()
    print("2. 转换HTML文件:")
    print("   python main.py -i examples/sample.html -o output.pptx")
    print()
    print("3. 查看帮助:")
    print("   python main.py --help")
    print()
    print("4. 如果遇到问题:")
    print("   查看 TROUBLESHOOTING.md 文件")

def main():
    """主函数"""
    print("🎯 HTML2PPT 快速开始")
    print("=" * 50)
    
    # 检测环境
    env_info = detect_environment()
    print_environment_info(env_info)
    
    # 推荐安装方式
    method, message = recommend_installation(env_info)
    
    if method == "error":
        print(f"❌ {message}")
        return 1
    
    print(f"💡 建议: {message}")
    
    # 询问用户是否继续
    try:
        response = input("\n是否继续安装? (y/N): ").strip().lower()
        if response not in ['y', 'yes', '是']:
            print("安装已取消")
            return 0
    except KeyboardInterrupt:
        print("\n安装已取消")
        return 0
    
    # 执行安装
    if run_installation(method, env_info):
        print("\n✅ 安装完成")
        
        # 测试安装
        if test_installation():
            show_next_steps()
            return 0
        else:
            print("\n⚠️  安装可能不完整，请查看错误信息")
            return 1
    else:
        print("\n❌ 安装失败")
        print("\n🔧 故障排除建议:")
        print("1. 检查网络连接")
        print("2. 尝试使用代理或镜像源")
        print("3. 查看 TROUBLESHOOTING.md")
        print("4. 手动安装: pip install -r requirements-minimal.txt")
        return 1

if __name__ == "__main__":
    sys.exit(main())