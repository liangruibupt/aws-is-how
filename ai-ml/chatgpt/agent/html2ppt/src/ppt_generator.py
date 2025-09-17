"""
PPT生成器 - 根据分析结果生成PowerPoint演示文稿
"""

import os
import logging
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from PIL import Image
import re

try:
    from .style_mapper import StyleMapper
    from .template_manager import TemplateManager
except ImportError:
    # Handle direct execution or testing
    from style_mapper import StyleMapper
    from template_manager import TemplateManager

logger = logging.getLogger(__name__)

class PPTGenerator:
    """PowerPoint生成器"""
    
    def __init__(self, config=None):
        self.config = config or {}
        self.ppt_config = self.config.get('ppt_generation', {})
        self.layout_config = self.config.get('layout_settings', {})
        
        self.presentation = None
        self.style_mapper = StyleMapper(config)
        self.template_manager = TemplateManager(config)
        
        # 动态字体大小配置
        self.font_sizes = self._calculate_adaptive_font_sizes()
        
        # 统计信息
        self.stats = {
            'slides': 0,
            'images': 0,
            'text_blocks': 0,
            'tables': 0
        }
    
    def _calculate_adaptive_font_sizes(self):
        """根据幻灯片尺寸计算自适应字体大小"""
        slide_width = self.layout_config.get('slide_width', 10)  # inches
        slide_height = self.layout_config.get('slide_height', 7.5)  # inches
        
        # 使用更合理的字体大小计算，避免文字超出幻灯片边界
        # 基于标准演示文稿的最佳实践
        slide_diagonal = (slide_width ** 2 + slide_height ** 2) ** 0.5
        
        # 基础字体大小：使用更保守的计算方式，确保文字适合幻灯片
        # 标准10x7.5英寸幻灯片的基础字体大小应该在18-24pt之间
        base_size = max(18, min(24, int(slide_diagonal * 1.8)))
        
        # 优化的字体大小层次，确保内容不会超出边界
        return {
            'title': base_size + 12,      # 30-36pt (大标题)
            'h1': base_size + 8,          # 26-32pt (主标题)
            'h2': base_size + 6,          # 24-30pt (副标题)
            'h3': base_size + 4,          # 22-28pt (三级标题)
            'h4': base_size + 2,          # 20-26pt (四级标题)
            'h5': base_size,              # 18-24pt (五级标题)
            'h6': base_size - 2,          # 16-22pt (六级标题)
            'body': base_size,            # 18-24pt (正文)
            'list': base_size,            # 18-24pt (列表)
            'code': base_size - 4,        # 14-20pt (代码)
            'caption': base_size - 6      # 12-18pt (说明文字)
        }
    
    def set_template(self, template_path):
        """设置PPT模版"""
        self.template_manager.load_template(template_path)
    
    def create_presentation(self, analyzed_content, output_path):
        """创建PowerPoint演示文稿"""
        try:
            # 初始化演示文稿
            self._initialize_presentation()
            
            # 获取PPT结构
            ppt_structure = analyzed_content.get('ppt_structure', {})
            slide_plan = ppt_structure.get('slide_plan', [])
            
            if not slide_plan:
                logger.warning("没有找到幻灯片计划，使用默认结构")
                slide_plan = self._create_default_slide_plan(analyzed_content)
            
            # 生成幻灯片
            for slide_info in slide_plan:
                self._create_slide(slide_info, analyzed_content)
            
            # 应用设计建议
            design_suggestions = ppt_structure.get('design_suggestions', {})
            self._apply_design_suggestions(design_suggestions)
            
            # 保存演示文稿
            self._save_presentation(output_path, analyzed_content)
            
            logger.info(f"成功生成PPT: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"生成PPT失败: {e}")
            return False
    
    def _initialize_presentation(self):
        """初始化演示文稿"""
        template_path = self.template_manager.get_template_path()
        
        if template_path and os.path.exists(template_path):
            self.presentation = Presentation(template_path)
            logger.info(f"使用模版: {template_path}")
        else:
            self.presentation = Presentation()
            logger.info("使用默认模版")
        
        # 设置幻灯片尺寸
        slide_width = self.layout_config.get('slide_width', 10)
        slide_height = self.layout_config.get('slide_height', 7.5)
        
        self.presentation.slide_width = Inches(slide_width)
        self.presentation.slide_height = Inches(slide_height)
    
    def _create_default_slide_plan(self, analyzed_content):
        """创建默认幻灯片计划"""
        slide_plan = []
        
        # 标题页
        doc_info = analyzed_content.get('document_info', {})
        slide_plan.append({
            'slide_number': 1,
            'type': 'title',
            'title': doc_info.get('title', '演示文稿'),
            'content': '基于HTML文档自动生成',
            'layout': 'title_slide'
        })
        
        # 内容页
        sections = analyzed_content.get('sections', [])
        for i, section in enumerate(sections):
            slide_plan.append({
                'slide_number': i + 2,
                'type': 'content',
                'title': section.get('title', f'章节 {i+1}'),
                'content': section.get('content', []),
                'layout': 'content_slide',
                'source_section': i
            })
        
        return slide_plan
    
    def _force_disable_bullet(self, paragraph):
        """强制禁用段落的项目符号，确保模板不会覆盖"""
        try:
            paragraph.bullet = False
            # 额外确保没有项目符号格式
            if hasattr(paragraph, '_element'):
                # 移除任何项目符号相关的XML元素
                pPr = paragraph._element.get_or_add_pPr()
                # 移除项目符号设置
                for buChar in pPr.xpath('.//a:buChar', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}):
                    buChar.getparent().remove(buChar)
                for buAutoNum in pPr.xpath('.//a:buAutoNum', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}):
                    buAutoNum.getparent().remove(buAutoNum)
        except:
            # 如果XML操作失败，至少确保基本设置
            paragraph.bullet = False
    
    def _force_enable_bullet(self, paragraph):
        """强制启用段落的项目符号"""
        paragraph.bullet = True
    
    def _force_disable_bullet_add_number(self, paragraph, number):
        """强制禁用项目符号并添加数字前缀"""
        self._force_disable_bullet(paragraph)
        # 数字前缀已经在文本中添加了
    
    def _calculate_optimal_font_size(self, formatted_text_parts, font_size_key='body'):
        """计算最优字体大小，防止文字超出幻灯片边界"""
        if not formatted_text_parts:
            return self.font_sizes[font_size_key]
        
        # 计算文本总长度
        total_text_length = sum(len(part.get('text', '')) for part in formatted_text_parts)
        
        # 获取基础字体大小
        base_size = self.font_sizes[font_size_key]
        
        # 根据文本长度和类型调整字体大小
        if font_size_key == 'title':
            # 标题字体大小调整更保守
            if total_text_length > 50:
                return max(base_size - 8, 24)
            elif total_text_length > 30:
                return max(base_size - 4, 28)
            else:
                return base_size
        elif font_size_key.startswith('h'):
            # 标题字体大小调整
            if total_text_length > 80:
                return max(base_size - 6, 16)
            elif total_text_length > 50:
                return max(base_size - 3, 18)
            else:
                return base_size
        else:
            # 正文字体大小调整
            if total_text_length > 300:  # 很长文本
                return max(base_size - 6, 12)  # 大幅减小字体
            elif total_text_length > 200:  # 长文本
                return max(base_size - 4, 14)  # 减小字体
            elif total_text_length > 100:  # 中等长度文本
                return max(base_size - 2, 16)  # 稍微减小字体
            else:
                return base_size  # 短文本使用标准字体大小
    
    def _create_slide(self, slide_info, analyzed_content):
        """创建单张幻灯片"""
        slide_type = slide_info.get('type', 'content')
        
        if slide_type == 'title':
            self._create_title_slide(slide_info)
        elif slide_type == 'content':
            self._create_content_slide(slide_info, analyzed_content)
        elif slide_type == 'section':
            self._create_section_slide(slide_info)
        
        self.stats['slides'] += 1
    
    def _create_title_slide(self, slide_info):
        """创建标题幻灯片"""
        # 使用标题布局
        slide_layout = self.presentation.slide_layouts[0]  # 通常是标题布局
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide_info.get('title', '演示文稿')
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._format_title_text(slide.shapes.title)
        
        # 设置副标题 - 更安全的方式
        subtitle = slide_info.get('content', '')
        subtitle_placeholder = None
        
        # 查找副标题占位符
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 4:  # 副标题占位符 (SUBTITLE)
                subtitle_placeholder = placeholder
                break
        
        if subtitle_placeholder and subtitle:
            subtitle_placeholder.text = subtitle
            self._format_subtitle_text(subtitle_placeholder)
        elif subtitle:
            # 如果没有副标题占位符，创建文本框
            self._add_subtitle_textbox(slide, subtitle)
        
        # 添加生成时间到页脚
        timestamp = datetime.now().strftime("%Y年%m月%d日")
        self._add_footer_text(slide, f"生成时间: {timestamp}")
    
    def _add_subtitle_textbox(self, slide, subtitle):
        """添加副标题文本框"""
        # 在标题下方添加副标题
        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(1)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = subtitle
        
        # 设置副标题样式，使用优化的字体大小
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        optimal_size = self._calculate_optimal_font_size([{'text': subtitle}], 'body')
        for run in paragraph.runs:
            run.font.size = Pt(optimal_size)
            run.font.color.rgb = RGBColor(0, 0, 0)  # 使用黑色确保可见性
    
    def _create_content_slide(self, slide_info, analyzed_content):
        """创建内容幻灯片"""
        content = slide_info.get('content', [])
        layout_type = self._determine_best_layout(content, slide_info.get('layout', 'content_slide'))
        
        # 选择合适的布局
        slide_layout = self._select_slide_layout(layout_type)
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide_info.get('title', '')
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._format_heading_text(slide.shapes.title)
        
        # 添加内容
        self._add_slide_content(slide, content, analyzed_content, title)
        
        # 添加备注
        self._add_slide_notes(slide, slide_info, analyzed_content)
    
    def _determine_best_layout(self, content, suggested_layout):
        """根据内容确定最佳布局"""
        if not content:
            return 'content_slide'
        
        # 分析内容类型
        has_table = any(item.get('type') == 'table' for item in content)
        has_image = any(item.get('type') == 'img' for item in content)
        has_list = any(item.get('type') in ['ul', 'ol'] for item in content)
        content_count = len(content)
        
        # 根据内容特征选择布局
        if has_table:
            return 'content_slide'  # 表格用标准内容布局
        elif has_image and content_count <= 2:
            return 'content_slide'  # 图片内容用标准布局
        elif has_list and content_count <= 3:
            return 'content_slide'  # 列表用标准布局
        elif content_count == 1:
            return 'content_slide'
        else:
            return suggested_layout
    
    def _create_section_slide(self, slide_info):
        """创建章节分隔幻灯片"""
        slide_layout = self.presentation.slide_layouts[2]  # 通常是章节标题布局
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title = slide_info.get('title', '章节')
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._format_section_title_text(slide.shapes.title)
    
    def _select_slide_layout(self, layout_type):
        """选择幻灯片布局"""
        layout_mapping = {
            'title_slide': 0,
            'content_slide': 1,
            'section_slide': 2,
            'image_focus': 6,
            'table_layout': 4,
            'bullet_points': 1,
            'single_content': 5,
            'multi_content': 1,
            'content_heavy': 1
        }
        
        layout_index = layout_mapping.get(layout_type, 1)
        
        # 确保布局索引有效
        if layout_index >= len(self.presentation.slide_layouts):
            layout_index = 1
        
        return self.presentation.slide_layouts[layout_index]
    
    def _add_slide_content(self, slide, content, analyzed_content, slide_title=""):
        """添加幻灯片内容"""
        if not content:
            return
        
        # 尝试使用内容占位符
        content_placeholder = self._find_content_placeholder(slide)
        
        if content_placeholder and self._can_use_placeholder(content):
            # 使用占位符添加内容
            self._add_content_to_placeholder(content_placeholder, content, slide_title)
        else:
            # 使用自定义布局添加内容
            self._add_content_with_custom_layout(slide, content, slide_title)
    
    def _find_content_placeholder(self, slide):
        """查找内容占位符"""
        for placeholder in slide.placeholders:
            # 内容占位符类型通常是2或7
            if placeholder.placeholder_format.type in [2, 7]:
                return placeholder
        return None
    
    def _can_use_placeholder(self, content):
        """判断是否可以使用占位符"""
        # 临时禁用占位符，强制使用自定义布局来解决项目符号问题
        return False
        
        # 原来的逻辑（暂时禁用）
        # text_types = ['p', 'div', 'ul', 'ol', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'pre', 'code', 'blockquote']
        # for item in content:
        #     if item.get('type') not in text_types:
        #         return False
        # return len(content) <= 5  # 内容不太多时使用占位符
    
    def _add_content_to_placeholder(self, placeholder, content, slide_title=""):
        """向占位符添加内容"""
        if not placeholder.has_text_frame:
            return
        
        text_frame = placeholder.text_frame
        text_frame.clear()
        
        # 添加内容到占位符
        paragraph_index = 0
        
        for i, item in enumerate(content):
            item_type = item.get('type', '')
            text = item.get('text', '').strip()
            
            # 跳过空内容，但不跳过列表（列表没有直接文本）
            if not text and item_type not in ['ul', 'ol']:
                continue
            
            # 跳过与幻灯片标题重复的标题元素
            if item_type in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6'] and text == slide_title:
                continue
            
            if item_type in ['ul', 'ol']:
                # 处理列表 - 区分有序和无序列表
                items = item.get('items', [])
                for j, list_item in enumerate(items):
                    list_text = list_item.get('text', '').strip()
                    formatted_text = list_item.get('formatted_text', [])
                    
                    if list_text or formatted_text:
                        if paragraph_index == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.level = 0
                        
                        # 根据列表类型设置不同的格式
                        if item_type == 'ol':
                            # 有序列表 - 手动添加数字，强制禁用项目符号
                            prefix = f"{j + 1}. "
                            self._force_disable_bullet_add_number(p, j + 1)
                        else:
                            # 无序列表 - 强制启用PowerPoint项目符号
                            prefix = ""
                            self._force_enable_bullet(p)
                        
                        # 使用格式化文本或纯文本
                        if formatted_text:
                            # 为有序列表添加数字前缀
                            if prefix:
                                prefix_part = [{'text': prefix, 'bold': False, 'italic': False, 'underline': False, 'link': None, 'color': None}]
                                combined_text = prefix_part + formatted_text
                                self._add_formatted_text_to_paragraph(p, combined_text, 'list')
                            else:
                                self._add_formatted_text_to_paragraph(p, formatted_text, 'list')
                        else:
                            p.text = prefix + list_text
                            # 设置列表样式，使用优化的字体大小
                            optimal_size = self._calculate_optimal_font_size([{'text': list_text}], 'list')
                            for run in p.runs:
                                run.font.size = Pt(optimal_size)
                        
                        paragraph_index += 1
            elif item_type in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                # 处理标题（已经过滤掉重复的）
                if paragraph_index == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.level = 0
                # 标题不应该有项目符号 - 强制设置
                self._force_disable_bullet(p)
                
                formatted_text = item.get('formatted_text', [])
                
                if formatted_text:
                    level = int(item_type[1])
                    self._add_formatted_text_to_paragraph(p, formatted_text, f'h{level}')
                    # 设置标题样式
                    for run in p.runs:
                        run.font.bold = True
                        run.font.size = Pt(self.font_sizes[f'h{level}'])
                else:
                    p.text = text
                    # 设置标题样式，使用优化的字体大小
                    level = int(item_type[1])
                    optimal_size = self._calculate_optimal_font_size([{'text': text}], f'h{level}')
                    for run in p.runs:
                        run.font.bold = True
                        run.font.size = Pt(optimal_size)
                
                paragraph_index += 1
            elif item_type in ['pre', 'code']:
                # 处理代码块
                if paragraph_index == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.level = 0
                # 代码块不应该有项目符号 - 强制设置
                self._force_disable_bullet(p)
                
                formatted_text = item.get('formatted_text', [])
                
                if formatted_text:
                    self._add_formatted_text_to_paragraph(p, formatted_text, 'code')
                else:
                    p.text = text
                
                # 设置代码样式，使用优化的字体大小
                optimal_size = self._calculate_optimal_font_size([{'text': text}], 'code')
                for run in p.runs:
                    run.font.name = 'Consolas'
                    run.font.size = Pt(optimal_size)
                    run.font.color.rgb = RGBColor(0, 0, 0)
                
                paragraph_index += 1
            elif item_type in ['blockquote']:
                # 处理引用
                if paragraph_index == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.level = 0
                # 引用不应该有项目符号 - 强制设置
                self._force_disable_bullet(p)
                
                formatted_text = item.get('formatted_text', [])
                
                if formatted_text:
                    # 为引用添加引号
                    quote_parts = [{'text': '"', 'bold': False, 'italic': True, 'underline': False, 'link': None, 'color': None}]
                    quote_parts.extend(formatted_text)
                    quote_parts.append({'text': '"', 'bold': False, 'italic': True, 'underline': False, 'link': None, 'color': None})
                    self._add_formatted_text_to_paragraph(p, quote_parts, 'body')
                else:
                    p.text = f'"{text}"'  # 添加引号
                
                # 设置引用样式，使用优化的字体大小
                optimal_size = self._calculate_optimal_font_size([{'text': text}], 'body')
                for run in p.runs:
                    run.font.italic = True
                    run.font.size = Pt(optimal_size)
                
                paragraph_index += 1
            else:
                # 处理普通段落
                if paragraph_index == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.level = 0
                # 普通段落不应该有项目符号 - 强制设置
                self._force_disable_bullet(p)
                
                formatted_text = item.get('formatted_text', [])
                
                if formatted_text:
                    self._add_formatted_text_to_paragraph(p, formatted_text, 'body')
                else:
                    p.text = text
                    # 设置普通段落样式，使用优化的字体大小
                    optimal_size = self._calculate_optimal_font_size([{'text': text}], 'body')
                    for run in p.runs:
                        run.font.size = Pt(optimal_size)
                
                paragraph_index += 1
            
            # 应用样式
            if paragraph_index > 0:  # 确保段落存在
                styles = item.get('styles', {})
                current_p = text_frame.paragraphs[paragraph_index - 1]
                self.style_mapper.apply_text_styles(current_p, styles)
    
    def _add_content_with_custom_layout(self, slide, content, slide_title=""):
        """使用自定义布局添加内容"""
        margins = self.layout_config.get('margins', {})
        slide_width = self.layout_config.get('slide_width', 10)
        slide_height = self.layout_config.get('slide_height', 7.5)
        
        # 计算可用区域
        left_margin = margins.get('left', 0.5)
        top_margin = margins.get('top', 1.5)  # 为标题留出空间
        right_margin = margins.get('right', 0.5)
        bottom_margin = margins.get('bottom', 0.5)
        
        available_width = slide_width - left_margin - right_margin
        available_height = slide_height - top_margin - bottom_margin
        
        current_y = top_margin
        
        for item in content:
            if current_y >= slide_height - bottom_margin:
                break  # 空间不足，停止添加
            
            item_type = item.get('type', '')
            item_text = item.get('text', '').strip()
            
            # 跳过与幻灯片标题重复的标题元素
            if item_type in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6'] and item_text == slide_title:
                continue
            
            if item_type in ['p', 'div']:
                height = self._add_paragraph_content_at_position(slide, item, left_margin, current_y, available_width)
                current_y += height + 0.1  # 添加间距
            elif item_type in ['ul', 'ol']:
                height = self._add_list_content_at_position(slide, item, left_margin, current_y, available_width)
                current_y += height + 0.1
            elif item_type == 'table':
                height = self._add_table_content_at_position(slide, item, left_margin, current_y, available_width)
                current_y += height + 0.2
            elif item_type == 'img':
                height = self._add_image_content_at_position(slide, item, left_margin, current_y, available_width)
                current_y += height + 0.2
            elif item_type in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                height = self._add_heading_content_at_position(slide, item, left_margin, current_y, available_width)
                current_y += height + 0.1
            elif item_type in ['blockquote', 'pre', 'code']:
                height = self._add_code_content_at_position(slide, item, left_margin, current_y, available_width)
                current_y += height + 0.1
    
    def _create_content_textbox(self, slide):
        """创建内容文本框"""
        margins = self.layout_config.get('margins', {})
        left = Inches(margins.get('left', 0.5))
        top = Inches(margins.get('top', 1.5))
        width = Inches(self.layout_config.get('slide_width', 10) - margins.get('left', 0.5) - margins.get('right', 0.5))
        height = Inches(self.layout_config.get('slide_height', 7.5) - margins.get('top', 1.5) - margins.get('bottom', 0.5))
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        return textbox
    
    def _add_paragraph_content_at_position(self, slide, item, left, top, width):
        """在指定位置添加段落内容"""
        text = item.get('text', '').strip()
        formatted_text = item.get('formatted_text', [])
        
        if not text and not formatted_text:
            return 0
        
        # 估算文本高度（更准确的估算）
        display_text = text or ''.join([part['text'] for part in formatted_text])
        lines = max(1, len(display_text) // 80)  # 假设每行80个字符
        text_height = max(0.3, lines * 0.25)  # 每行约0.25英寸
        
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), 
            Inches(width), Inches(text_height)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        
        # 添加格式化文本或纯文本
        if formatted_text:
            # 先应用基础样式，再添加格式化文本（这样格式化文本的颜色会覆盖基础样式）
            text_frame.text = ""  # 确保有一个段落
            # 强制禁用段落的项目符号
            self._force_disable_bullet(text_frame.paragraphs[0])
            styles = item.get('styles', {})
            self.style_mapper.apply_text_styles(text_frame.paragraphs[0], styles)
            # 然后添加格式化文本，个别颜色会覆盖基础样式
            self._add_formatted_text_to_paragraph(text_frame.paragraphs[0], formatted_text, 'body')
        else:
            text_frame.text = text
            # 强制禁用段落的项目符号
            self._force_disable_bullet(text_frame.paragraphs[0])
            # 应用样式
            styles = item.get('styles', {})
            self.style_mapper.apply_text_styles(text_frame.paragraphs[0], styles)
            
            # 只有在样式中没有颜色时才设置默认黑色
            if not styles.get('color'):
                for run in text_frame.paragraphs[0].runs:
                    run.font.color.rgb = RGBColor(0, 0, 0)
        
        self.stats['text_blocks'] += 1
        return text_height
    
    def _add_formatted_text_to_paragraph(self, paragraph, formatted_text_parts, font_size_key='body'):
        """向段落添加格式化文本"""
        if not formatted_text_parts:
            return
        
        # 清空段落现有内容
        paragraph.clear()
        
        # 计算合适的字体大小，防止文字超出边界
        base_font_size = self._calculate_optimal_font_size(formatted_text_parts, font_size_key)
        
        for i, part in enumerate(formatted_text_parts):
            text = part.get('text', '')
            if not text:
                continue
            
            # 添加文本运行
            if i == 0:
                # 第一个部分使用段落的默认运行
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                run.text = text
            else:
                # 后续部分添加新的运行
                run = paragraph.add_run()
                run.text = text
            
            # 设置优化后的字体大小
            run.font.size = Pt(base_font_size)
            
            # 应用格式
            if part.get('bold', False):
                run.font.bold = True
            if part.get('italic', False):
                run.font.italic = True
            if part.get('underline', False):
                run.font.underline = True
            
            # 处理颜色
            color = part.get('color')
            if color:
                try:
                    # 解析颜色值
                    if color.startswith('#'):
                        # 十六进制颜色
                        hex_color = color[1:]
                        if len(hex_color) == 6:
                            r = int(hex_color[0:2], 16)
                            g = int(hex_color[2:4], 16)
                            b = int(hex_color[4:6], 16)
                            run.font.color.rgb = RGBColor(r, g, b)
                    elif color.startswith('rgb('):
                        # RGB颜色格式 rgb(255, 0, 0)
                        rgb_values = color[4:-1].split(',')
                        if len(rgb_values) == 3:
                            r = int(rgb_values[0].strip())
                            g = int(rgb_values[1].strip())
                            b = int(rgb_values[2].strip())
                            run.font.color.rgb = RGBColor(r, g, b)
                except:
                    # 如果颜色解析失败，使用黑色作为fallback
                    run.font.color.rgb = RGBColor(0, 0, 0)
            else:
                # 只有在没有指定颜色时才使用黑色作为默认
                run.font.color.rgb = RGBColor(0, 0, 0)
            
            # 处理链接（PowerPoint中链接处理比较复杂，这里简化处理）
            link = part.get('link')
            if link:
                # 为链接设置蓝色和下划线
                run.font.color.rgb = RGBColor(0, 0, 255)
                run.font.underline = True
    
    def _add_list_content_at_position(self, slide, item, left, top, width):
        """在指定位置添加列表内容"""
        items = item.get('items', [])
        item_type = item.get('type', 'ul')
        if not items:
            return 0
        
        # 估算列表高度
        list_height = max(0.5, len(items) * 0.25)
        
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), 
            Inches(width), Inches(list_height)
        )
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        
        # 添加列表项
        for i, list_item in enumerate(items):
            item_text = list_item.get('text', '').strip()
            formatted_text = list_item.get('formatted_text', [])
            
            if item_text or formatted_text:
                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                p.level = 0
                
                # 根据列表类型设置不同的格式
                if item_type == 'ol':
                    # 有序列表 - 手动添加数字，强制禁用项目符号
                    prefix = f"{i + 1}. "
                    self._force_disable_bullet_add_number(p, i + 1)
                else:
                    # 无序列表 - 强制启用PowerPoint项目符号
                    prefix = ""
                    self._force_enable_bullet(p)
                
                # 添加格式化文本或纯文本
                if formatted_text:
                    # 为有序列表添加数字前缀
                    if prefix:
                        prefix_part = [{'text': prefix, 'bold': False, 'italic': False, 'underline': False, 'link': None, 'color': None}]
                        combined_text = prefix_part + formatted_text
                        self._add_formatted_text_to_paragraph(p, combined_text, 'list')
                    else:
                        self._add_formatted_text_to_paragraph(p, formatted_text, 'list')
                else:
                    p.text = prefix + item_text
                    # 设置列表样式
                    for run in p.runs:
                        run.font.size = Pt(self.font_sizes['list'])
                
                # 应用自定义样式
                styles = list_item.get('styles', {})
                self.style_mapper.apply_text_styles(p, styles)
        
        self.stats['text_blocks'] += 1
        return list_height
    
    def _add_table_content(self, slide, item, y_offset):
        """添加表格内容"""
        table_data = item.get('table_data', {})
        if not table_data:
            return 0
        
        headers = table_data.get('headers', [])
        rows = table_data.get('rows', [])
        
        if not rows:
            return 0
        
        # 计算表格尺寸
        row_count = len(rows) + (1 if headers else 0)
        col_count = max(len(row) for row in rows) if rows else 1
        
        margins = self.layout_config.get('margins', {})
        left = Inches(margins.get('left', 0.5))
        top = Inches(margins.get('top', 1.5) + y_offset)
        width = Inches(self.layout_config.get('slide_width', 10) - margins.get('left', 0.5) - margins.get('right', 0.5))
        height = Inches(min(4.0, row_count * 0.4))
        
        # 创建表格
        table_shape = slide.shapes.add_table(row_count, col_count, left, top, width, height)
        table = table_shape.table
        
        # 添加表头
        current_row = 0
        if headers:
            for col, header in enumerate(headers):
                if col < col_count:
                    cell = table.cell(current_row, col)
                    cell.text = header.get('text', '')
                    # 表头样式
                    self._format_table_header_cell(cell)
            current_row += 1
        
        # 添加数据行
        for row_data in rows:
            if current_row < row_count:
                for col, cell_data in enumerate(row_data):
                    if col < col_count:
                        cell = table.cell(current_row, col)
                        cell.text = cell_data.get('text', '')
                        # 数据单元格样式
                        self._format_table_data_cell(cell)
                current_row += 1
        
        self.stats['tables'] += 1
        return height.inches
    
    def _add_image_content(self, slide, item, y_offset):
        """添加图片内容"""
        image_info = item.get('image_info', {})
        if not image_info:
            return 0
        
        local_path = image_info.get('local_path')
        if not local_path or not os.path.exists(local_path):
            # 如果图片不存在，添加占位符
            return self._add_image_placeholder(slide, image_info, y_offset)
        
        try:
            # 获取图片信息
            with Image.open(local_path) as img:
                img_width, img_height = img.size
            
            # 计算合适的显示尺寸
            max_width = self.layout_config.get('slide_width', 10) * 0.8
            max_height = (self.layout_config.get('slide_height', 7.5) - 2) * 0.6
            
            # 保持宽高比
            aspect_ratio = img_width / img_height
            if max_width / aspect_ratio <= max_height:
                display_width = max_width
                display_height = max_width / aspect_ratio
            else:
                display_height = max_height
                display_width = max_height * aspect_ratio
            
            # 居中放置
            margins = self.layout_config.get('margins', {})
            left = Inches((self.layout_config.get('slide_width', 10) - display_width) / 2)
            top = Inches(margins.get('top', 1.5) + y_offset)
            
            # 添加图片
            slide.shapes.add_picture(local_path, left, top, Inches(display_width), Inches(display_height))
            
            # 添加图片说明
            alt_text = image_info.get('alt', '') or image_info.get('title', '')
            if alt_text:
                caption_top = top + Inches(display_height) + Inches(0.1)
                caption_textbox = slide.shapes.add_textbox(
                    left, caption_top, Inches(display_width), Inches(0.3)
                )
                caption_textbox.text_frame.text = alt_text
                caption_textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # 设置说明文字样式
                for run in caption_textbox.text_frame.paragraphs[0].runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = RGBColor(128, 128, 128)
            
            self.stats['images'] += 1
            return display_height + 0.5  # 包括说明文字的高度
            
        except Exception as e:
            logger.error(f"添加图片失败 {local_path}: {e}")
            return self._add_image_placeholder(slide, image_info, y_offset)
    
    def _add_image_placeholder(self, slide, image_info, y_offset):
        """添加图片占位符"""
        margins = self.layout_config.get('margins', {})
        left = Inches(margins.get('left', 0.5))
        top = Inches(margins.get('top', 1.5) + y_offset)
        width = Inches(4)
        height = Inches(2)
        
        # 创建占位符形状
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        
        # 设置占位符样式
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
        placeholder.line.color.rgb = RGBColor(200, 200, 200)
        
        # 添加占位符文字
        placeholder.text_frame.text = f"图片: {image_info.get('alt', '无法加载')}"
        placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return 2.5
    
    def _add_heading_content(self, slide, item, y_offset):
        """添加标题内容"""
        text = item.get('text', '').strip()
        if not text:
            return 0
        
        margins = self.layout_config.get('margins', {})
        left = Inches(margins.get('left', 0.5))
        top = Inches(margins.get('top', 1.5) + y_offset)
        width = Inches(self.layout_config.get('slide_width', 10) - margins.get('left', 0.5) - margins.get('right', 0.5))
        height = Inches(0.8)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text
        
        # 根据标题级别设置样式
        heading_level = int(item.get('type', 'h3')[1])  # 提取数字
        self._format_heading_by_level(text_frame.paragraphs[0], heading_level)
        
        # 应用自定义样式
        styles = item.get('styles', {})
        self.style_mapper.apply_text_styles(text_frame.paragraphs[0], styles)
        
        self.stats['text_blocks'] += 1
        return 0.8
    
    def _add_code_content(self, slide, item, y_offset):
        """添加代码内容"""
        text = item.get('text', '').strip()
        if not text:
            return 0
        
        margins = self.layout_config.get('margins', {})
        left = Inches(margins.get('left', 0.5))
        top = Inches(margins.get('top', 1.5) + y_offset)
        width = Inches(self.layout_config.get('slide_width', 10) - margins.get('left', 0.5) - margins.get('right', 0.5))
        
        # 估算代码块高度
        lines = text.split('\n')
        code_height = max(0.5, len(lines) * 0.2)
        height = Inches(code_height)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text
        
        # 代码样式
        paragraph = text_frame.paragraphs[0]
        for run in paragraph.runs:
            run.font.name = 'Consolas'
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(0, 0, 0)
        
        # 设置背景色
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(248, 248, 248)
        textbox.line.color.rgb = RGBColor(200, 200, 200)
        
        self.stats['text_blocks'] += 1
        return code_height
    
    def _format_title_text(self, title_shape):
        """格式化标题文本"""
        if not title_shape.has_text_frame:
            return
        
        paragraph = title_shape.text_frame.paragraphs[0]
        for run in paragraph.runs:
            run.font.size = Pt(self.font_sizes['title'])
            run.font.bold = True
            run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    
    def _format_subtitle_text(self, subtitle_shape):
        """格式化副标题文本"""
        if not subtitle_shape.has_text_frame:
            return
        
        paragraph = subtitle_shape.text_frame.paragraphs[0]
        for run in paragraph.runs:
            run.font.size = Pt(self.font_sizes['body'])
            run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2
    
    def _format_heading_text(self, heading_shape):
        """格式化标题文本"""
        if not heading_shape.has_text_frame:
            return
        
        paragraph = heading_shape.text_frame.paragraphs[0]
        for run in paragraph.runs:
            run.font.size = Pt(self.font_sizes['h2'])
            run.font.bold = True
            run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    
    def _format_section_title_text(self, title_shape):
        """格式化章节标题文本"""
        if not title_shape.has_text_frame:
            return
        
        paragraph = title_shape.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(self.font_sizes['h1'])
            run.font.bold = True
            run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    
    def _format_heading_by_level(self, paragraph, level):
        """根据级别格式化标题"""
        size_mapping = {1: 28, 2: 24, 3: 20, 4: 18, 5: 16, 6: 14}
        font_size = size_mapping.get(level, 16)
        
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = True
            if level <= 2:
                run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
            else:
                run.font.color.theme_color = MSO_THEME_COLOR.DARK_1
    
    def _format_table_header_cell(self, cell):
        """格式化表头单元格"""
        cell.fill.solid()
        cell.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.font.size = Pt(self.font_sizes['body'] - 4)
    
    def _format_table_data_cell(self, cell):
        """格式化数据单元格"""
        paragraph = cell.text_frame.paragraphs[0]
        for run in paragraph.runs:
            run.font.size = Pt(self.font_sizes['body'] - 2)
            run.font.color.theme_color = MSO_THEME_COLOR.DARK_1
    
    def _add_footer_text(self, slide, text):
        """添加页脚文本"""
        margins = self.layout_config.get('margins', {})
        left = Inches(margins.get('left', 0.5))
        top = Inches(self.layout_config.get('slide_height', 7.5) - 0.5)
        width = Inches(self.layout_config.get('slide_width', 10) - margins.get('left', 0.5) - margins.get('right', 0.5))
        height = Inches(0.3)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text
        
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.RIGHT
        for run in paragraph.runs:
            run.font.size = Pt(self.font_sizes['caption'])
            run.font.color.rgb = RGBColor(128, 128, 128)
    
    def _add_slide_notes(self, slide, slide_info, analyzed_content):
        """添加幻灯片备注"""
        if not self.ppt_config.get('notes', {}).get('include_original_html', True):
            return
        
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        
        notes_content = []
        
        # 添加原始HTML信息
        if self.ppt_config.get('notes', {}).get('include_original_html', True):
            source_section = slide_info.get('source_section')
            if source_section is not None:
                sections = analyzed_content.get('sections', [])
                if source_section < len(sections):
                    section = sections[source_section]
                    notes_content.append(f"原始章节: {section.get('title', '')}")
                    
                    # 添加关键点
                    key_points = section.get('key_points', [])
                    if key_points:
                        notes_content.append("关键点:")
                        for point in key_points[:3]:  # 只添加前3个关键点
                            notes_content.append(f"- {point.get('text', '')}")
        
        # 添加元数据
        if self.ppt_config.get('notes', {}).get('include_metadata', True):
            notes_content.append(f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
            notes_content.append(f"幻灯片类型: {slide_info.get('type', 'content')}")
            notes_content.append(f"布局: {slide_info.get('layout', 'default')}")
        
        # 设置备注文本
        if notes_content:
            max_length = self.ppt_config.get('notes', {}).get('max_note_length', 1000)
            notes_text = '\n'.join(notes_content)[:max_length]
            notes_text_frame.text = notes_text
    
    def _apply_design_suggestions(self, design_suggestions):
        """应用设计建议"""
        if not design_suggestions:
            return
        
        try:
            # 应用配色方案
            color_scheme = design_suggestions.get('color_scheme', {})
            if color_scheme:
                self._apply_color_scheme(color_scheme)
            
            # 应用字体建议
            font_suggestions = design_suggestions.get('font_suggestions', {})
            if font_suggestions:
                self._apply_font_suggestions(font_suggestions)
                
        except Exception as e:
            logger.error(f"应用设计建议失败: {e}")
    
    def _apply_color_scheme(self, color_scheme):
        """应用配色方案"""
        # 这里可以实现更复杂的配色逻辑
        # 由于python-pptx的限制，主要通过模版来控制配色
        pass
    
    def _apply_font_suggestions(self, font_suggestions):
        """应用字体建议"""
        # 这里可以实现字体应用逻辑
        # 由于python-pptx的限制，主要在创建内容时应用字体
        pass
    
    def _save_presentation(self, output_path, analyzed_content):
        """保存演示文稿"""
        # 确保输出目录存在
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        
        # 设置文档属性
        doc_props = self.presentation.core_properties
        doc_props.author = self.ppt_config.get('metadata', {}).get('author', 'HTML2PPT Converter')
        doc_props.title = analyzed_content.get('document_info', {}).get('title', 'Converted Presentation')
        doc_props.subject = self.ppt_config.get('metadata', {}).get('subject', 'Auto-generated presentation')
        doc_props.created = datetime.now()
        doc_props.modified = datetime.now()
        
        # 保存文件
        self.presentation.save(output_path)
        
        logger.info(f"演示文稿已保存: {output_path}")
        logger.info(f"统计信息: {self.stats}")
    
    def get_statistics(self):
        """获取生成统计信息"""
        return self.stats.copy()
    
    def _add_heading_content_at_position(self, slide, item, left, top, width):
        """在指定位置添加标题内容"""
        text = item.get('text', '').strip()
        formatted_text = item.get('formatted_text', [])
        
        if not text and not formatted_text:
            return 0
        
        height = 0.6  # 标题固定高度
        
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), 
            Inches(width), Inches(height)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        # 添加格式化文本或纯文本
        paragraph = text_frame.paragraphs[0]
        heading_level = int(item.get('type', 'h3')[1])
        if formatted_text:
            self._add_formatted_text_to_paragraph(paragraph, formatted_text, f'h{heading_level}')
        else:
            text_frame.text = text
        
        # 根据标题级别设置样式
        self._format_heading_by_level(paragraph, heading_level)
        
        # 应用自定义样式
        styles = item.get('styles', {})
        self.style_mapper.apply_text_styles(paragraph, styles)
        
        self.stats['text_blocks'] += 1
        return height
    
    def _add_table_content_at_position(self, slide, item, left, top, width):
        """在指定位置添加表格内容"""
        table_data = item.get('table_data', {})
        if not table_data:
            return 0
        
        headers = table_data.get('headers', [])
        rows = table_data.get('rows', [])
        
        if not rows:
            return 0
        
        # 计算表格尺寸
        row_count = len(rows) + (1 if headers else 0)
        col_count = max(len(row) for row in rows) if rows else 1
        
        height = min(3.0, row_count * 0.4)  # 限制表格高度
        
        # 创建表格
        table_shape = slide.shapes.add_table(
            row_count, col_count, 
            Inches(left), Inches(top), 
            Inches(width), Inches(height)
        )
        table = table_shape.table
        
        # 添加表头
        current_row = 0
        if headers:
            for col, header in enumerate(headers):
                if col < col_count:
                    cell = table.cell(current_row, col)
                    cell.text = header.get('text', '')
                    self._format_table_header_cell(cell)
            current_row += 1
        
        # 添加数据行
        for row_data in rows:
            if current_row < row_count:
                for col, cell_data in enumerate(row_data):
                    if col < col_count:
                        cell = table.cell(current_row, col)
                        cell.text = cell_data.get('text', '')
                        self._format_table_data_cell(cell)
                current_row += 1
        
        self.stats['tables'] += 1
        return height
    
    def _add_image_content_at_position(self, slide, item, left, top, width):
        """在指定位置添加图片内容"""
        image_info = item.get('image_info', {})
        if not image_info:
            return 0
        
        local_path = image_info.get('local_path')
        if not local_path or not os.path.exists(local_path):
            return self._add_image_placeholder_at_position(slide, image_info, left, top, width)
        
        try:
            # 获取图片信息并计算尺寸
            with Image.open(local_path) as img:
                img_width, img_height = img.size
            
            # 计算合适的显示尺寸
            max_width = width * 0.8
            max_height = 3.0  # 限制图片高度
            
            # 保持宽高比
            aspect_ratio = img_width / img_height
            if max_width / aspect_ratio <= max_height:
                display_width = max_width
                display_height = max_width / aspect_ratio
            else:
                display_height = max_height
                display_width = max_height * aspect_ratio
            
            # 居中放置，添加顶部间距
            image_left = left + (width - display_width) / 2
            image_top = top + 0.1  # 添加0.1英寸的顶部间距
            
            # 添加图片
            slide.shapes.add_picture(
                local_path, 
                Inches(image_left), Inches(image_top), 
                Inches(display_width), Inches(display_height)
            )
            
            # 添加图片说明
            alt_text = image_info.get('alt', '') or image_info.get('title', '')
            caption_height = 0
            if alt_text:
                caption_top = image_top + display_height + 0.1  # 使用调整后的image_top
                caption_textbox = slide.shapes.add_textbox(
                    Inches(image_left), Inches(caption_top), 
                    Inches(display_width), Inches(0.3)
                )
                caption_textbox.text_frame.text = alt_text
                caption_textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # 设置说明文字样式
                for run in caption_textbox.text_frame.paragraphs[0].runs:
                    run.font.size = Pt(self.font_sizes['caption'])
                    run.font.color.rgb = RGBColor(128, 128, 128)
                
                caption_height = 0.4  # 说明文字高度
            
            self.stats['images'] += 1
            return display_height + 0.1 + caption_height + 0.1  # 顶部间距 + 图片高度 + 说明文字 + 底部间距
            
        except Exception as e:
            logger.error(f"添加图片失败 {local_path}: {e}")
            return self._add_image_placeholder_at_position(slide, image_info, left, top, width)
    
    def _add_image_placeholder_at_position(self, slide, image_info, left, top, width):
        """在指定位置添加图片占位符"""
        height = 2.0
        placeholder_width = min(width, 4.0)
        placeholder_left = left + (width - placeholder_width) / 2
        placeholder_top = top + 0.1  # 添加顶部间距
        
        # 创建占位符形状
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(placeholder_left), Inches(placeholder_top), 
            Inches(placeholder_width), Inches(height)
        )
        
        # 设置占位符样式
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
        placeholder.line.color.rgb = RGBColor(200, 200, 200)
        
        # 添加占位符文字
        placeholder.text_frame.text = f"图片: {image_info.get('alt', '无法加载')}"
        placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return height + 0.1 + 0.1  # 顶部间距 + 占位符高度 + 底部间距
    
    def _add_code_content_at_position(self, slide, item, left, top, width):
        """在指定位置添加代码内容"""
        text = item.get('text', '').strip()
        formatted_text = item.get('formatted_text', [])
        
        if not text and not formatted_text:
            return 0
        
        # 估算代码块高度
        display_text = text or ''.join([part['text'] for part in formatted_text])
        lines = display_text.split('\n')
        code_height = max(0.5, len(lines) * 0.2)
        
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), 
            Inches(width), Inches(code_height)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        # 添加格式化文本或纯文本
        paragraph = text_frame.paragraphs[0]
        if formatted_text:
            self._add_formatted_text_to_paragraph(paragraph, formatted_text, 'code')
        else:
            text_frame.text = text
        
        # 代码样式
        for run in paragraph.runs:
            run.font.name = 'Consolas'
            run.font.size = Pt(self.font_sizes['code'])
            run.font.color.rgb = RGBColor(0, 0, 0)
        
        # 设置背景色
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(248, 248, 248)
        textbox.line.color.rgb = RGBColor(200, 200, 200)
        
        self.stats['text_blocks'] += 1
        return code_height
    
    def _format_title_text(self, title_shape):
        """格式化标题文本"""
        if not title_shape.has_text_frame:
            return
        
        text_frame = title_shape.text_frame
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                # 使用优化的标题字体大小
                optimal_size = self._calculate_optimal_font_size([{'text': run.text}], 'title')
                run.font.size = Pt(optimal_size)
                run.font.bold = True
                # 让样式映射器处理颜色，或使用深蓝色作为标题默认色
                if not hasattr(run.font.color, 'rgb') or run.font.color.rgb is None:
                    run.font.color.rgb = RGBColor(0, 120, 212)  # 深蓝色 #0078D4
    
    def _format_subtitle_text(self, subtitle_shape):
        """格式化副标题文本"""
        if not subtitle_shape.has_text_frame:
            return
        
        text_frame = subtitle_shape.text_frame
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                # 使用优化的副标题字体大小
                optimal_size = self._calculate_optimal_font_size([{'text': run.text}], 'body')
                run.font.size = Pt(optimal_size)
                # 使用深灰色作为副标题默认色
                if not hasattr(run.font.color, 'rgb') or run.font.color.rgb is None:
                    run.font.color.rgb = RGBColor(51, 51, 51)  # 深灰色 #333
    
    def _format_heading_text(self, heading_shape):
        """格式化标题文本"""
        if not heading_shape.has_text_frame:
            return
        
        text_frame = heading_shape.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                # 使用优化的标题字体大小
                optimal_size = self._calculate_optimal_font_size([{'text': run.text}], 'h1')
                run.font.size = Pt(optimal_size)
                run.font.bold = True
                # 使用深蓝色作为标题默认色
                if not hasattr(run.font.color, 'rgb') or run.font.color.rgb is None:
                    run.font.color.rgb = RGBColor(0, 120, 212)  # 深蓝色 #0078D4
    
    def _format_section_title_text(self, title_shape):
        """格式化章节标题文本"""
        if not title_shape.has_text_frame:
            return
        
        text_frame = title_shape.text_frame
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                # 使用优化的章节标题字体大小
                optimal_size = self._calculate_optimal_font_size([{'text': run.text}], 'h1')
                run.font.size = Pt(optimal_size)
                run.font.bold = True
                # 使用深蓝色作为章节标题默认色
                if not hasattr(run.font.color, 'rgb') or run.font.color.rgb is None:
                    run.font.color.rgb = RGBColor(0, 120, 212)  # 深蓝色 #0078D4
    
    def _add_footer_text(self, slide, footer_text):
        """添加页脚文本"""
        # 在幻灯片底部添加页脚
        left = Inches(0.5)
        top = Inches(self.layout_config.get('slide_height', 7.5) - 0.5)
        width = Inches(self.layout_config.get('slide_width', 10) - 1)
        height = Inches(0.3)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = footer_text
        
        # 设置页脚样式
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.RIGHT
        for run in paragraph.runs:
            run.font.size = Pt(self.font_sizes['caption'])
            run.font.color.rgb = RGBColor(128, 128, 128)  # 使用灰色作为页脚
    
    def _add_slide_notes(self, slide, slide_info, analyzed_content):
        """添加幻灯片备注"""
        # 可以在这里添加演讲者备注
        notes_text = slide_info.get('notes', '')
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text