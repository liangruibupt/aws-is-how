"""
样式映射器 - 将HTML样式映射到PowerPoint样式
"""

import logging
import re
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR

logger = logging.getLogger(__name__)

class StyleMapper:
    """样式映射器"""
    
    def __init__(self, config=None):
        self.config = config or {}
        self.style_config = self.config.get('style_mapping', {})
        
        # 默认样式映射
        self.tag_mapping = self.style_config.get('tag_mapping', {
            'h1': 'title',
            'h2': 'subtitle',
            'h3': 'heading',
            'h4': 'subheading',
            'p': 'body',
            'ul': 'bullet_list',
            'ol': 'number_list',
            'blockquote': 'quote',
            'code': 'code'
        })
        
        self.font_mapping = self.style_config.get('font_mapping', {
            'default': 'Microsoft YaHei',
            'serif': 'Times New Roman',
            'sans-serif': 'Arial',
            'monospace': 'Consolas'
        })
        
        self.color_mapping = self.style_config.get('color_mapping', {
            'preserve_original': True,
            'fallback_colors': {
                'text': '#000000',
                'background': '#FFFFFF',
                'accent': '#0078D4'
            }
        })
    
    def apply_text_styles(self, paragraph, html_styles):
        """应用文本样式到段落"""
        if not html_styles:
            return
        
        try:
            # 应用字体样式
            self._apply_font_styles(paragraph, html_styles)
            
            # 应用颜色样式
            self._apply_color_styles(paragraph, html_styles)
            
            # 应用对齐样式
            self._apply_alignment_styles(paragraph, html_styles)
            
            # 应用间距样式
            self._apply_spacing_styles(paragraph, html_styles)
            
        except Exception as e:
            logger.error(f"应用文本样式失败: {e}")
    
    def _apply_font_styles(self, paragraph, html_styles):
        """应用字体样式"""
        for run in paragraph.runs:
            # 保存现有格式状态
            existing_bold = run.font.bold
            existing_italic = run.font.italic
            existing_underline = run.font.underline
            
            # 字体族
            font_family = html_styles.get('font-family', '').strip('"\'')
            if font_family:
                mapped_font = self._map_font_family(font_family)
                if mapped_font:
                    run.font.name = mapped_font
            
            # 字体大小
            font_size = html_styles.get('font-size', '')
            if font_size:
                size_pt = self._parse_font_size(font_size)
                if size_pt:
                    run.font.size = Pt(size_pt)
            
            # 字体粗细 - 只在明确指定时才覆盖现有格式
            font_weight = html_styles.get('font-weight', '')
            if font_weight in ['bold', 'bolder', '700', '800', '900']:
                run.font.bold = True
            elif font_weight in ['normal', '400'] and font_weight:
                # 只有在明确设置为normal时才重置
                run.font.bold = False
            # 如果没有font-weight样式，保持现有的bold状态
            
            # 字体样式 - 只在明确指定时才覆盖现有格式
            font_style = html_styles.get('font-style', '')
            if font_style == 'italic':
                run.font.italic = True
            elif font_style == 'normal' and font_style:
                # 只有在明确设置为normal时才重置
                run.font.italic = False
            # 如果没有font-style样式，保持现有的italic状态
            
            # 文本装饰
            text_decoration = html_styles.get('text-decoration', '')
            if 'underline' in text_decoration:
                run.font.underline = True
    
    def _apply_color_styles(self, paragraph, html_styles):
        """应用颜色样式"""
        # 文本颜色
        color = html_styles.get('color', '')
        if color:
            rgb_color = self._parse_color(color)
            if rgb_color:
                for run in paragraph.runs:
                    run.font.color.rgb = rgb_color
        
        # 背景颜色（通过段落级别设置）
        background_color = html_styles.get('background-color', '')
        if background_color:
            # PowerPoint段落背景色设置比较复杂，这里暂时跳过
            pass
    
    def _apply_alignment_styles(self, paragraph, html_styles):
        """应用对齐样式"""
        text_align = html_styles.get('text-align', '')
        
        alignment_mapping = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        
        if text_align in alignment_mapping:
            paragraph.alignment = alignment_mapping[text_align]
    
    def _apply_spacing_styles(self, paragraph, html_styles):
        """应用间距样式"""
        # 行高
        line_height = html_styles.get('line-height', '')
        if line_height:
            try:
                if line_height.endswith('%'):
                    # 百分比行高
                    percent = float(line_height[:-1]) / 100
                    paragraph.line_spacing = percent
                elif line_height.replace('.', '').isdigit():
                    # 数字倍数
                    multiplier = float(line_height)
                    paragraph.line_spacing = multiplier
            except ValueError:
                pass
        
        # 段落间距
        margin_top = html_styles.get('margin-top', '')
        margin_bottom = html_styles.get('margin-bottom', '')
        
        if margin_top:
            space_before = self._parse_spacing(margin_top)
            if space_before:
                paragraph.space_before = Pt(space_before)
        
        if margin_bottom:
            space_after = self._parse_spacing(margin_bottom)
            if space_after:
                paragraph.space_after = Pt(space_after)
    
    def _map_font_family(self, font_family):
        """映射字体族"""
        # 清理字体名称
        font_family = font_family.lower().strip()
        
        # 直接映射
        if font_family in self.font_mapping:
            return self.font_mapping[font_family]
        
        # 模糊匹配
        if 'serif' in font_family:
            return self.font_mapping.get('serif', 'Times New Roman')
        elif 'sans' in font_family or 'arial' in font_family:
            return self.font_mapping.get('sans-serif', 'Arial')
        elif 'mono' in font_family or 'courier' in font_family or 'consolas' in font_family:
            return self.font_mapping.get('monospace', 'Consolas')
        elif any(chinese_font in font_family for chinese_font in ['微软雅黑', 'yahei', '宋体', 'simsun', '黑体', 'simhei']):
            return 'Microsoft YaHei'
        
        # 默认字体
        return self.font_mapping.get('default', 'Microsoft YaHei')
    
    def _parse_font_size(self, font_size):
        """解析字体大小"""
        try:
            if font_size.endswith('px'):
                # 像素转点
                px = float(font_size[:-2])
                return px * 0.75  # 1px ≈ 0.75pt
            elif font_size.endswith('pt'):
                # 点
                return float(font_size[:-2])
            elif font_size.endswith('em'):
                # em转点（假设基础字体为12pt）
                em = float(font_size[:-2])
                return em * 12
            elif font_size.endswith('%'):
                # 百分比转点（假设基础字体为12pt）
                percent = float(font_size[:-1])
                return (percent / 100) * 12
            elif font_size.replace('.', '').isdigit():
                # 纯数字，假设为像素
                px = float(font_size)
                return px * 0.75
            else:
                # 关键词映射
                size_keywords = {
                    'xx-small': 8,
                    'x-small': 9,
                    'small': 10,
                    'medium': 12,
                    'large': 14,
                    'x-large': 18,
                    'xx-large': 24
                }
                return size_keywords.get(font_size.lower(), 12)
        except ValueError:
            return None
    
    def _parse_color(self, color_value):
        """解析颜色值"""
        try:
            color_value = color_value.strip().lower()
            
            # 十六进制颜色
            if color_value.startswith('#'):
                hex_color = color_value[1:]
                if len(hex_color) == 3:
                    # 短格式 #RGB -> #RRGGBB
                    hex_color = ''.join([c*2 for c in hex_color])
                
                if len(hex_color) == 6:
                    r = int(hex_color[0:2], 16)
                    g = int(hex_color[2:4], 16)
                    b = int(hex_color[4:6], 16)
                    return RGBColor(r, g, b)
            
            # RGB颜色
            elif color_value.startswith('rgb('):
                rgb_match = re.match(r'rgb\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)', color_value)
                if rgb_match:
                    r, g, b = map(int, rgb_match.groups())
                    return RGBColor(r, g, b)
            
            # RGBA颜色（忽略透明度）
            elif color_value.startswith('rgba('):
                rgba_match = re.match(r'rgba\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*,\s*[\d.]+\s*\)', color_value)
                if rgba_match:
                    r, g, b = map(int, rgba_match.groups())
                    return RGBColor(r, g, b)
            
            # 颜色名称
            else:
                color_names = {
                    'black': RGBColor(0, 0, 0),
                    'white': RGBColor(255, 255, 255),
                    'red': RGBColor(255, 0, 0),
                    'green': RGBColor(0, 128, 0),
                    'blue': RGBColor(0, 0, 255),
                    'yellow': RGBColor(255, 255, 0),
                    'cyan': RGBColor(0, 255, 255),
                    'magenta': RGBColor(255, 0, 255),
                    'gray': RGBColor(128, 128, 128),
                    'grey': RGBColor(128, 128, 128),
                    'orange': RGBColor(255, 165, 0),
                    'purple': RGBColor(128, 0, 128),
                    'brown': RGBColor(165, 42, 42),
                    'pink': RGBColor(255, 192, 203),
                    'lime': RGBColor(0, 255, 0),
                    'navy': RGBColor(0, 0, 128),
                    'teal': RGBColor(0, 128, 128),
                    'silver': RGBColor(192, 192, 192),
                    'maroon': RGBColor(128, 0, 0),
                    'olive': RGBColor(128, 128, 0)
                }
                
                return color_names.get(color_value)
        
        except (ValueError, AttributeError):
            pass
        
        return None
    
    def _parse_spacing(self, spacing_value):
        """解析间距值"""
        try:
            spacing_value = spacing_value.strip().lower()
            
            if spacing_value.endswith('px'):
                # 像素转点
                px = float(spacing_value[:-2])
                return px * 0.75
            elif spacing_value.endswith('pt'):
                # 点
                return float(spacing_value[:-2])
            elif spacing_value.endswith('em'):
                # em转点（假设基础字体为12pt）
                em = float(spacing_value[:-2])
                return em * 12
            elif spacing_value.replace('.', '').isdigit():
                # 纯数字，假设为像素
                px = float(spacing_value)
                return px * 0.75
        
        except ValueError:
            pass
        
        return None
    
    def get_ppt_style_for_html_tag(self, tag_name):
        """获取HTML标签对应的PPT样式"""
        return self.tag_mapping.get(tag_name.lower(), 'body')
    
    def apply_list_styles(self, text_frame, list_type='ul', html_styles=None):
        """应用列表样式"""
        if not text_frame.paragraphs:
            return
        
        # 设置列表缩进
        for paragraph in text_frame.paragraphs:
            if list_type == 'ol':
                # 有序列表
                paragraph.level = 0
            else:
                # 无序列表
                paragraph.level = 0
            
            # 应用HTML样式
            if html_styles:
                self.apply_text_styles(paragraph, html_styles)
    
    def apply_table_styles(self, table, html_styles=None):
        """应用表格样式"""
        if not table or not html_styles:
            return
        
        try:
            # 表格边框
            border_style = html_styles.get('border', '')
            if border_style:
                # 这里可以实现表格边框样式
                pass
            
            # 表格背景色
            background_color = html_styles.get('background-color', '')
            if background_color:
                rgb_color = self._parse_color(background_color)
                if rgb_color:
                    # 应用到所有单元格
                    for row in table.rows:
                        for cell in row.cells:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = rgb_color
        
        except Exception as e:
            logger.error(f"应用表格样式失败: {e}")
    
    def create_style_summary(self, html_styles):
        """创建样式摘要"""
        summary = {}
        
        if not html_styles:
            return summary
        
        # 提取关键样式信息
        if 'font-family' in html_styles:
            summary['font'] = self._map_font_family(html_styles['font-family'])
        
        if 'font-size' in html_styles:
            summary['size'] = self._parse_font_size(html_styles['font-size'])
        
        if 'color' in html_styles:
            summary['color'] = self._parse_color(html_styles['color'])
        
        if 'font-weight' in html_styles:
            summary['bold'] = html_styles['font-weight'] in ['bold', 'bolder', '700', '800', '900']
        
        if 'font-style' in html_styles:
            summary['italic'] = html_styles['font-style'] == 'italic'
        
        if 'text-align' in html_styles:
            summary['alignment'] = html_styles['text-align']
        
        return summary