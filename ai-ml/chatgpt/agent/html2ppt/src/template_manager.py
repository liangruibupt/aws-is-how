"""
模版管理器 - 管理PowerPoint模版
"""

import os
import logging
import shutil
from pathlib import Path
from pptx import Presentation

logger = logging.getLogger(__name__)

class TemplateManager:
    """模版管理器"""
    
    def __init__(self, config=None):
        self.config = config or {}
        self.ppt_config = self.config.get('ppt_generation', {})
        
        self.template_path = None
        self.template_info = {}
        
        # 默认模版路径
        self.default_template = self.ppt_config.get('default_template', 'templates/default.pptx')
        
        # 初始化模版目录
        self._initialize_template_directory()
    
    def _initialize_template_directory(self):
        """初始化模版目录"""
        template_dir = Path('templates')
        template_dir.mkdir(exist_ok=True)
        
        # 如果默认模版不存在，创建一个基础模版
        default_template_path = Path(self.default_template)
        if not default_template_path.exists():
            self._create_default_template(default_template_path)
    
    def _create_default_template(self, template_path):
        """创建默认模版"""
        try:
            # 创建基础演示文稿
            prs = Presentation()
            
            # 设置幻灯片尺寸
            slide_width = self.config.get('layout_settings', {}).get('slide_width', 10)
            slide_height = self.config.get('layout_settings', {}).get('slide_height', 7.5)
            
            from pptx.util import Inches
            prs.slide_width = Inches(slide_width)
            prs.slide_height = Inches(slide_height)
            
            # 确保目录存在
            template_path.parent.mkdir(parents=True, exist_ok=True)
            
            # 保存默认模版
            prs.save(str(template_path))
            
            logger.info(f"创建默认模版: {template_path}")
            
        except Exception as e:
            logger.error(f"创建默认模版失败: {e}")
    
    def load_template(self, template_path):
        """加载模版"""
        try:
            if not os.path.exists(template_path):
                logger.error(f"模版文件不存在: {template_path}")
                return False
            
            # 验证模版文件
            if not self._validate_template(template_path):
                logger.error(f"模版文件无效: {template_path}")
                return False
            
            self.template_path = template_path
            self.template_info = self._analyze_template(template_path)
            
            logger.info(f"成功加载模版: {template_path}")
            return True
            
        except Exception as e:
            logger.error(f"加载模版失败: {e}")
            return False
    
    def _validate_template(self, template_path):
        """验证模版文件"""
        try:
            # 尝试打开PowerPoint文件
            prs = Presentation(template_path)
            
            # 检查是否有布局
            if len(prs.slide_layouts) == 0:
                logger.warning("模版没有幻灯片布局")
                return False
            
            return True
            
        except Exception as e:
            logger.error(f"模版验证失败: {e}")
            return False
    
    def _analyze_template(self, template_path):
        """分析模版信息"""
        info = {
            'path': template_path,
            'layouts': [],
            'master_slides': [],
            'color_scheme': {},
            'font_scheme': {}
        }
        
        try:
            prs = Presentation(template_path)
            
            # 分析布局
            for i, layout in enumerate(prs.slide_layouts):
                layout_info = {
                    'index': i,
                    'name': layout.name,
                    'placeholders': []
                }
                
                # 分析占位符
                for placeholder in layout.placeholders:
                    placeholder_info = {
                        'index': placeholder.placeholder_format.idx,
                        'type': placeholder.placeholder_format.type,
                        'name': getattr(placeholder, 'name', f'Placeholder {placeholder.placeholder_format.idx}')
                    }
                    layout_info['placeholders'].append(placeholder_info)
                
                info['layouts'].append(layout_info)
            
            # 分析母版
            for i, master in enumerate(prs.slide_masters):
                master_info = {
                    'index': i,
                    'name': getattr(master, 'name', f'Master {i}')
                }
                info['master_slides'].append(master_info)
            
            logger.info(f"模版分析完成: {len(info['layouts'])} 个布局, {len(info['master_slides'])} 个母版")
            
        except Exception as e:
            logger.error(f"模版分析失败: {e}")
        
        return info
    
    def get_template_path(self):
        """获取当前模版路径"""
        if self.template_path and os.path.exists(self.template_path):
            return self.template_path
        
        # 返回默认模版
        default_path = Path(self.default_template)
        if default_path.exists():
            return str(default_path)
        
        return None
    
    def get_template_info(self):
        """获取模版信息"""
        return self.template_info.copy()
    
    def list_available_templates(self):
        """列出可用模版"""
        templates = []
        template_dir = Path('templates')
        
        if template_dir.exists():
            for template_file in template_dir.glob('*.pptx'):
                if self._validate_template(str(template_file)):
                    templates.append({
                        'name': template_file.stem,
                        'path': str(template_file),
                        'size': template_file.stat().st_size
                    })
        
        return templates
    
    def install_template(self, source_path, template_name=None):
        """安装新模版"""
        try:
            source_path = Path(source_path)
            
            if not source_path.exists():
                logger.error(f"源模版文件不存在: {source_path}")
                return False
            
            # 验证源模版
            if not self._validate_template(str(source_path)):
                logger.error(f"源模版文件无效: {source_path}")
                return False
            
            # 确定目标路径
            if not template_name:
                template_name = source_path.stem
            
            target_path = Path('templates') / f"{template_name}.pptx"
            
            # 复制模版文件
            shutil.copy2(str(source_path), str(target_path))
            
            logger.info(f"模版安装成功: {target_path}")
            return True
            
        except Exception as e:
            logger.error(f"安装模版失败: {e}")
            return False
    
    def remove_template(self, template_name):
        """删除模版"""
        try:
            template_path = Path('templates') / f"{template_name}.pptx"
            
            if not template_path.exists():
                logger.error(f"模版不存在: {template_path}")
                return False
            
            # 不允许删除默认模版
            if str(template_path) == self.default_template:
                logger.error("不能删除默认模版")
                return False
            
            template_path.unlink()
            
            logger.info(f"模版删除成功: {template_path}")
            return True
            
        except Exception as e:
            logger.error(f"删除模版失败: {e}")
            return False
    
    def get_layout_by_type(self, layout_type):
        """根据类型获取布局索引"""
        layout_mapping = {
            'title': 0,
            'content': 1,
            'section': 2,
            'comparison': 3,
            'content_with_caption': 4,
            'blank': 5,
            'picture_with_caption': 6
        }
        
        # 从模版信息中查找
        if self.template_info and 'layouts' in self.template_info:
            layouts = self.template_info['layouts']
            
            # 尝试按名称匹配
            for layout in layouts:
                layout_name = layout.get('name', '').lower()
                if layout_type.lower() in layout_name:
                    return layout['index']
        
        # 使用默认映射
        return layout_mapping.get(layout_type, 1)
    
    def get_suitable_layout(self, content_info):
        """根据内容信息获取合适的布局"""
        content_types = content_info.get('content_types', [])
        content_count = content_info.get('content_count', 0)
        
        # 根据内容特征选择布局
        if 'image' in content_types and content_count <= 2:
            return self.get_layout_by_type('picture_with_caption')
        elif 'table' in content_types:
            return self.get_layout_by_type('content')
        elif content_count == 0:
            return self.get_layout_by_type('blank')
        elif content_count == 1:
            return self.get_layout_by_type('content_with_caption')
        else:
            return self.get_layout_by_type('content')
    
    def create_custom_template(self, template_name, design_config):
        """创建自定义模版"""
        try:
            from pptx.util import Inches
            from pptx.dml.color import RGBColor
            from pptx.enum.dml import MSO_THEME_COLOR
            
            # 创建新演示文稿
            prs = Presentation()
            
            # 设置幻灯片尺寸
            slide_width = design_config.get('slide_width', 10)
            slide_height = design_config.get('slide_height', 7.5)
            prs.slide_width = Inches(slide_width)
            prs.slide_height = Inches(slide_height)
            
            # 应用配色方案
            color_scheme = design_config.get('color_scheme', {})
            if color_scheme:
                self._apply_color_scheme_to_template(prs, color_scheme)
            
            # 应用字体方案
            font_scheme = design_config.get('font_scheme', {})
            if font_scheme:
                self._apply_font_scheme_to_template(prs, font_scheme)
            
            # 保存自定义模版
            template_path = Path('templates') / f"{template_name}.pptx"
            template_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(template_path))
            
            logger.info(f"自定义模版创建成功: {template_path}")
            return str(template_path)
            
        except Exception as e:
            logger.error(f"创建自定义模版失败: {e}")
            return None
    
    def _apply_color_scheme_to_template(self, presentation, color_scheme):
        """应用配色方案到模版"""
        # 这里可以实现更复杂的配色逻辑
        # 由于python-pptx的限制，主要通过母版来控制配色
        pass
    
    def _apply_font_scheme_to_template(self, presentation, font_scheme):
        """应用字体方案到模版"""
        # 这里可以实现字体方案应用逻辑
        # 由于python-pptx的限制，主要通过母版来控制字体
        pass
    
    def backup_template(self, template_name):
        """备份模版"""
        try:
            template_path = Path('templates') / f"{template_name}.pptx"
            
            if not template_path.exists():
                logger.error(f"模版不存在: {template_path}")
                return False
            
            # 创建备份目录
            backup_dir = Path('templates') / 'backups'
            backup_dir.mkdir(exist_ok=True)
            
            # 生成备份文件名
            from datetime import datetime
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            backup_path = backup_dir / f"{template_name}_{timestamp}.pptx"
            
            # 复制文件
            shutil.copy2(str(template_path), str(backup_path))
            
            logger.info(f"模版备份成功: {backup_path}")
            return str(backup_path)
            
        except Exception as e:
            logger.error(f"备份模版失败: {e}")
            return None
    
    def restore_template(self, backup_path, template_name):
        """恢复模版"""
        try:
            backup_path = Path(backup_path)
            
            if not backup_path.exists():
                logger.error(f"备份文件不存在: {backup_path}")
                return False
            
            # 验证备份文件
            if not self._validate_template(str(backup_path)):
                logger.error(f"备份文件无效: {backup_path}")
                return False
            
            # 目标路径
            target_path = Path('templates') / f"{template_name}.pptx"
            
            # 复制文件
            shutil.copy2(str(backup_path), str(target_path))
            
            logger.info(f"模版恢复成功: {target_path}")
            return True
            
        except Exception as e:
            logger.error(f"恢复模版失败: {e}")
            return False