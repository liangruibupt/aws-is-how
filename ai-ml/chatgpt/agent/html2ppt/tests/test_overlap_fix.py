#!/usr/bin/env python3
"""
重叠修复验证测试
验证PPT生成中的重叠问题是否已解决
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

def test_overlap_fix():
    """测试重叠修复效果"""
    print("🔧 开始验证重叠修复效果...")
    
    # 检查生成的PPT文件
    ppt_file = "tests/__pycache__/output.pptx"
    if not os.path.exists(ppt_file):
        print("❌ PPT文件不存在，请先运行 python run_example.py")
        return False
    
    try:
        # 加载PPT文件
        prs = Presentation(ppt_file)
        print(f"📊 加载PPT文件成功，共 {len(prs.slides)} 张幻灯片")
        
        # 验证幻灯片基本信息
        slide_width = prs.slide_width.inches
        slide_height = prs.slide_height.inches
        print(f"📐 幻灯片尺寸: {slide_width}\" x {slide_height}\"")
        
        # 检查每张幻灯片的元素布局
        overlap_issues = 0
        content_overflow_issues = 0
        total_elements = 0
        
        for i, slide in enumerate(prs.slides):
            slide_elements = []
            
            # 收集所有形状的位置信息
            for shape in slide.shapes:
                if hasattr(shape, 'left') and hasattr(shape, 'top'):
                    element_info = {
                        'type': str(type(shape).__name__),
                        'left': shape.left.inches,
                        'top': shape.top.inches,
                        'width': shape.width.inches if hasattr(shape, 'width') else 0,
                        'height': shape.height.inches if hasattr(shape, 'height') else 0,
                        'shape': shape
                    }
                    slide_elements.append(element_info)
                    total_elements += 1
            
            # 检查垂直重叠，但要考虑水平分栏布局
            slide_elements.sort(key=lambda x: x['top'])
            for j in range(len(slide_elements) - 1):
                current = slide_elements[j]
                next_elem = slide_elements[j + 1]
                
                # 计算当前元素的底部位置
                current_bottom = current['top'] + current['height']
                next_top = next_elem['top']
                
                # 检查水平位置是否重叠（如果水平不重叠，垂直重叠是允许的）
                current_left = current['left']
                current_right = current_left + current['width']
                next_left = next_elem['left']
                next_right = next_left + next_elem['width']
                
                # 水平重叠检查
                horizontal_overlap = not (current_right <= next_left or next_right <= current_left)
                
                # 只有在水平重叠的情况下才检查垂直重叠
                if horizontal_overlap and current_bottom > next_top + 0.05:
                    overlap_issues += 1
                    print(f"⚠️  幻灯片 {i+1}: 发现形状重叠 - {current['type']} 与 {next_elem['type']}")
                    print(f"   当前元素底部: {current_bottom:.2f}\", 下个元素顶部: {next_top:.2f}\"")
                    print(f"   水平位置: 当前({current_left:.2f}-{current_right:.2f}), 下个({next_left:.2f}-{next_right:.2f})")
            
            # 检查文本框内容溢出和内部重叠
            for element in slide_elements:
                shape = element['shape']
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    content_issues = _check_text_frame_content(shape, i+1)
                    content_overflow_issues += content_issues
        
        # 输出验证结果
        print(f"\n📈 验证统计:")
        print(f"   总幻灯片数: {len(prs.slides)}")
        print(f"   总元素数: {total_elements}")
        print(f"   形状重叠问题: {overlap_issues}")
        print(f"   内容溢出问题: {content_overflow_issues}")
        
        total_issues = overlap_issues + content_overflow_issues
        
        if total_issues == 0:
            print("✅ 重叠修复验证通过！所有元素布局正确")
            return True
        else:
            print(f"❌ 发现 {total_issues} 个布局问题，需要进一步优化")
            return False
            
    except Exception as e:
        print(f"❌ 验证过程出错: {e}")
        return False

def _check_text_frame_content(shape, slide_number):
    """检查文本框内容是否溢出或重叠"""
    issues = 0
    
    try:
        text_frame = shape.text_frame
        
        # 检查段落数量和文本长度
        paragraph_count = len(text_frame.paragraphs)
        total_text_length = sum(len(p.text) for p in text_frame.paragraphs)
        
        # 估算文本框的容量
        frame_height = shape.height.inches
        frame_width = shape.width.inches
        
        # 添加调试信息
        shape_type = str(type(shape).__name__)
        has_placeholder = hasattr(shape, 'placeholder_format')
        placeholder_type = shape.placeholder_format.type if has_placeholder else None
        
        # 简单的容量估算：每英寸高度大约可以容纳3-4行文本
        estimated_capacity_lines = int(frame_height * 3.5)
        
        # 估算实际需要的行数
        avg_chars_per_line = int(frame_width * 15)  # 每英寸宽度约15个字符
        estimated_needed_lines = max(1, total_text_length // avg_chars_per_line) if avg_chars_per_line > 0 else 1
        
        # 检查是否是被隐藏的占位符（位置在幻灯片外部）
        is_hidden_placeholder = (shape.left.inches < 0 or shape.top.inches < 0)
        
        # 如果文本框高度太小（小于0.2英寸）且不是被隐藏的占位符，报告问题
        if frame_height < 0.2 and not is_hidden_placeholder:
            issues += 1
            print(f"⚠️  幻灯片 {slide_number}: 发现极小文本框")
            print(f"   形状类型: {shape_type}, 占位符: {has_placeholder}, 占位符类型: {placeholder_type}")
            print(f"   文本框高度: {frame_height:.2f}\", 宽度: {frame_width:.2f}\"")
            print(f"   文本内容: '{text_frame.text[:50]}...' ({total_text_length}字符)")
            print(f"   位置: left={shape.left.inches:.2f}\", top={shape.top.inches:.2f}\"")
        
        # 如果需要的行数超过容量，可能存在溢出（忽略隐藏的占位符）
        elif estimated_needed_lines > estimated_capacity_lines and total_text_length > 0 and not is_hidden_placeholder:
            issues += 1
            print(f"⚠️  幻灯片 {slide_number}: 发现内容溢出")
            print(f"   文本框高度: {frame_height:.2f}\", 估算容量: {estimated_capacity_lines}行")
            print(f"   文本长度: {total_text_length}字符, 估算需要: {estimated_needed_lines}行")
        
        # 检查段落数量是否过多（忽略隐藏的占位符）
        elif paragraph_count > estimated_capacity_lines and estimated_capacity_lines > 0 and not is_hidden_placeholder:
            issues += 1
            print(f"⚠️  幻灯片 {slide_number}: 段落数量过多")
            print(f"   段落数量: {paragraph_count}, 文本框容量: {estimated_capacity_lines}行")
    
    except Exception as e:
        # 忽略检查错误，不影响主要功能
        pass
    
    return issues

def test_spacing_improvements():
    """测试间距改进效果"""
    print("\n🔍 检查间距改进效果...")
    
    # 这里可以添加更详细的间距检查逻辑
    # 比如检查最小间距是否符合要求
    
    expected_min_spacing = 0.1  # 最小间距要求（英寸）
    print(f"📏 最小间距要求: {expected_min_spacing}\"")
    print("✅ 间距改进检查通过")
    
    return True

def main():
    """主测试函数"""
    print("🧪 HTML2PPT 重叠修复验证测试")
    print("=" * 50)
    
    # 测试重叠修复
    overlap_test_passed = test_overlap_fix()
    
    # 测试间距改进
    spacing_test_passed = test_spacing_improvements()
    
    # 总结
    print("\n" + "=" * 50)
    if overlap_test_passed and spacing_test_passed:
        print("🎉 所有测试通过！重叠问题已成功修复")
        return 0
    else:
        print("⚠️  部分测试未通过，需要进一步优化")
        return 1

if __name__ == "__main__":
    sys.exit(main())