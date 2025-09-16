"""
PPT生成器测试
"""

import unittest
import os
import sys
import tempfile
from pathlib import Path

# 添加src目录到Python路径
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

from ppt_generator import PPTGenerator

class TestPPTGenerator(unittest.TestCase):
    """PPT生成器测试类"""
    
    def setUp(self):
        """测试前准备"""
        self.config = {
            'layout_settings': {
                'slide_width': 10,
                'slide_height': 7.5,
                'margins': {
                    'top': 0.5,
                    'bottom': 0.5,
                    'left': 0.5,
                    'right': 0.5
                }
            },
            'ppt_generation': {
                'slide_types': {
                    'title_slide': True,
                    'content_slides': True,
                    'section_dividers': True
                },
                'notes': {
                    'include_original_html': True,
                    'include_metadata': True,
                    'max_note_length': 1000
                }
            }
        }
        
        self.generator = PPTGenerator(self.config)
        
        # 测试数据
        self.test_analyzed_content = {
            'document_info': {
                'title': '测试演示文稿',
                'meta': {'description': '这是一个测试演示文稿'},
                'language': 'zh-CN'
            },
            'sections': [
                {
                    'type': 'section',
                    'level': 1,
                    'title': '第一章节',
                    'content': [
                        {
                            'type': 'p',
                            'text': '这是第一个段落内容。',
                            'styles': {'color': '#000000', 'font-size': '14px'}
                        },
                        {
                            'type': 'ul',
                            'items': [
                                {'text': '列表项目1', 'styles': {}},
                                {'text': '列表项目2', 'styles': {}}
                            ]
                        }
                    ],
                    'importance_score': 0.8,
                    'content_types': ['paragraph', 'list'],
                    'suggested_layout': 'bullet_points'
                }
            ],
            'ppt_structure': {
                'total_slides': 2,
                'slide_plan': [
                    {
                        'slide_number': 1,
                        'type': 'title',
                        'title': '测试演示文稿',
                        'content': '基于HTML文档自动生成',
                        'layout': 'title_slide'
                    },
                    {
                        'slide_number': 2,
                        'type': 'content',
                        'title': '第一章节',
                        'content': [
                            {
                                'type': 'p',
                                'text': '这是第一个段落内容。',
                                'styles': {'color': '#000000', 'font-size': '14px'}
                            },
                            {
                                'type': 'ul',
                                'items': [
                                    {'text': '列表项目1', 'styles': {}},
                                    {'text': '列表项目2', 'styles': {}}
                                ]
                            }
                        ],
                        'layout': 'bullet_points',
                        'source_section': 0
                    }
                ],
                'design_suggestions': {
                    'color_scheme': {
                        'primary': '#0078D4',
                        'secondary': '#106EBE',
                        'accent': '#005A9E'
                    },
                    'font_suggestions': {
                        'title': 'Microsoft YaHei',
                        'heading': 'Microsoft YaHei',
                        'body': 'Microsoft YaHei'
                    }
                }
            }
        }
    
    def test_initialize_presentation(self):
        """测试演示文稿初始化"""
        self.generator._initialize_presentation()
        
        self.assertIsNotNone(self.generator.presentation)
        
        # 检查幻灯片尺寸
        from pptx.util import Inches
        expected_width = Inches(10)
        expected_height = Inches(7.5)
        
        self.assertEqual(self.generator.presentation.slide_width, expected_width)
        self.assertEqual(self.generator.presentation.slide_height, expected_height)
    
    def test_create_default_slide_plan(self):
        """测试默认幻灯片计划创建"""
        slide_plan = self.generator._create_default_slide_plan(self.test_analyzed_content)
        
        self.assertIsInstance(slide_plan, list)
        self.assertGreater(len(slide_plan), 0)
        
        # 检查标题页
        title_slide = slide_plan[0]
        self.assertEqual(title_slide['type'], 'title')
        self.assertEqual(title_slide['slide_number'], 1)
        self.assertEqual(title_slide['title'], '测试演示文稿')
    
    def test_select_slide_layout(self):
        """测试幻灯片布局选择"""
        self.generator._initialize_presentation()
        
        # 测试不同布局类型
        layout_types = ['title_slide', 'content_slide', 'bullet_points', 'image_focus']
        
        for layout_type in layout_types:
            layout = self.generator._select_slide_layout(layout_type)
            self.assertIsNotNone(layout)
    
    def test_create_slide(self):
        """测试幻灯片创建"""
        self.generator._initialize_presentation()
        
        # 测试标题幻灯片创建
        title_slide_info = {
            'type': 'title',
            'title': '测试标题',
            'content': '测试副标题'
        }
        
        initial_slide_count = len(self.generator.presentation.slides)
        self.generator._create_slide(title_slide_info, self.test_analyzed_content)
        
        # 检查是否添加了幻灯片
        self.assertEqual(len(self.generator.presentation.slides), initial_slide_count + 1)
        self.assertEqual(self.generator.stats['slides'], 1)
    
    def test_create_title_slide(self):
        """测试标题幻灯片创建"""
        self.generator._initialize_presentation()
        
        # Ensure we start with a clean presentation
        self.assertEqual(len(self.generator.presentation.slides), 0)
        
        slide_info = {
            'title': '测试演示文稿',
            'content': '自动生成的演示文稿'
        }
        
        self.generator._create_title_slide(slide_info)
        
        # 检查幻灯片是否创建
        self.assertEqual(len(self.generator.presentation.slides), 1)
        
        slide = self.generator.presentation.slides[0]
        self.assertIsNotNone(slide.shapes.title)
        self.assertEqual(slide.shapes.title.text, '测试演示文稿')
    
    def test_create_content_slide(self):
        """测试内容幻灯片创建"""
        self.generator._initialize_presentation()
        
        slide_info = {
            'title': '内容标题',
            'content': [
                {
                    'type': 'p',
                    'text': '这是段落内容',
                    'styles': {}
                }
            ],
            'layout': 'content_slide'
        }
        
        self.generator._create_content_slide(slide_info, self.test_analyzed_content)
        
        # 检查幻灯片是否创建
        self.assertEqual(len(self.generator.presentation.slides), 1)
        
        slide = self.generator.presentation.slides[0]
        self.assertIsNotNone(slide.shapes.title)
        self.assertEqual(slide.shapes.title.text, '内容标题')
    
    def test_add_paragraph_content(self):
        """测试段落内容添加"""
        self.generator._initialize_presentation()
        
        # 创建测试幻灯片
        slide_layout = self.generator.presentation.slide_layouts[1]
        slide = self.generator.presentation.slides.add_slide(slide_layout)
        
        item = {
            'type': 'p',
            'text': '这是一个测试段落，包含一些文本内容。',
            'styles': {'color': '#000000', 'font-size': '14px'}
        }
        
        height = self.generator._add_paragraph_content_at_position(slide, item, 0.5, 1.5, 8)
        
        self.assertGreater(height, 0)
        self.assertEqual(self.generator.stats['text_blocks'], 1)
    
    def test_add_list_content(self):
        """测试列表内容添加"""
        self.generator._initialize_presentation()
        
        # 创建测试幻灯片
        slide_layout = self.generator.presentation.slide_layouts[1]
        slide = self.generator.presentation.slides.add_slide(slide_layout)
        
        item = {
            'type': 'ul',
            'items': [
                {'text': '第一个列表项', 'styles': {}},
                {'text': '第二个列表项', 'styles': {}},
                {'text': '第三个列表项', 'styles': {}}
            ]
        }
        
        height = self.generator._add_list_content_at_position(slide, item, 0.5, 1.5, 8)
        
        self.assertGreater(height, 0)
        self.assertEqual(self.generator.stats['text_blocks'], 1)
    
    def test_add_table_content(self):
        """测试表格内容添加"""
        self.generator._initialize_presentation()
        
        # 创建测试幻灯片
        slide_layout = self.generator.presentation.slide_layouts[1]
        slide = self.generator.presentation.slides.add_slide(slide_layout)
        
        item = {
            'type': 'table',
            'table_data': {
                'headers': [
                    {'text': '列1', 'styles': {}},
                    {'text': '列2', 'styles': {}}
                ],
                'rows': [
                    [
                        {'text': '数据1', 'styles': {}},
                        {'text': '数据2', 'styles': {}}
                    ],
                    [
                        {'text': '数据3', 'styles': {}},
                        {'text': '数据4', 'styles': {}}
                    ]
                ]
            }
        }
        
        height = self.generator._add_table_content(slide, item, 0)
        
        self.assertGreater(height, 0)
        self.assertEqual(self.generator.stats['tables'], 1)
    
    def test_add_heading_content(self):
        """测试标题内容添加"""
        self.generator._initialize_presentation()
        
        # 创建测试幻灯片
        slide_layout = self.generator.presentation.slide_layouts[1]
        slide = self.generator.presentation.slides.add_slide(slide_layout)
        
        item = {
            'type': 'h2',
            'text': '这是一个二级标题',
            'styles': {'color': '#0078D4', 'font-weight': 'bold'}
        }
        
        height = self.generator._add_heading_content(slide, item, 0)
        
        self.assertGreater(height, 0)
        self.assertEqual(self.generator.stats['text_blocks'], 1)
    
    def test_add_code_content(self):
        """测试代码内容添加"""
        self.generator._initialize_presentation()
        
        # 创建测试幻灯片
        slide_layout = self.generator.presentation.slide_layouts[1]
        slide = self.generator.presentation.slides.add_slide(slide_layout)
        
        item = {
            'type': 'code',
            'text': 'def hello_world():\n    print("Hello, World!")\n    return True',
            'styles': {'font-family': 'monospace'}
        }
        
        height = self.generator._add_code_content(slide, item, 0)
        
        self.assertGreater(height, 0)
        self.assertEqual(self.generator.stats['text_blocks'], 1)
    
    def test_add_image_placeholder(self):
        """测试图片占位符添加"""
        self.generator._initialize_presentation()
        
        # 创建测试幻灯片
        slide_layout = self.generator.presentation.slide_layouts[1]
        slide = self.generator.presentation.slides.add_slide(slide_layout)
        
        image_info = {
            'src': 'nonexistent.jpg',
            'alt': '测试图片',
            'title': '这是一个测试图片'
        }
        
        height = self.generator._add_image_placeholder(slide, image_info, 0)
        
        self.assertGreater(height, 0)
    
    def test_add_slide_notes(self):
        """测试幻灯片备注添加"""
        self.generator._initialize_presentation()
        
        # 创建测试幻灯片
        slide_layout = self.generator.presentation.slide_layouts[1]
        slide = self.generator.presentation.slides.add_slide(slide_layout)
        
        slide_info = {
            'type': 'content',
            'title': '测试幻灯片',
            'source_section': 0
        }
        
        self.generator._add_slide_notes(slide, slide_info, self.test_analyzed_content)
        
        # 检查备注是否添加
        notes_slide = slide.notes_slide
        self.assertIsNotNone(notes_slide.notes_text_frame.text)
    
    def test_create_presentation_with_temp_file(self):
        """测试完整演示文稿创建（使用临时文件）"""
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            temp_path = temp_file.name
        
        try:
            success = self.generator.create_presentation(self.test_analyzed_content, temp_path)
            
            self.assertTrue(success)
            self.assertTrue(os.path.exists(temp_path))
            
            # 检查统计信息
            stats = self.generator.get_statistics()
            self.assertGreater(stats['slides'], 0)
            
        finally:
            # 清理临时文件
            if os.path.exists(temp_path):
                os.unlink(temp_path)
    
    def test_get_statistics(self):
        """测试统计信息获取"""
        stats = self.generator.get_statistics()
        
        self.assertIsInstance(stats, dict)
        self.assertIn('slides', stats)
        self.assertIn('images', stats)
        self.assertIn('text_blocks', stats)
        self.assertIn('tables', stats)
        
        # 初始统计应该都是0
        for value in stats.values():
            self.assertIsInstance(value, int)
            self.assertGreaterEqual(value, 0)
    
    def test_format_heading_by_level(self):
        """测试按级别格式化标题"""
        self.generator._initialize_presentation()
        
        # 创建测试幻灯片和段落
        slide_layout = self.generator.presentation.slide_layouts[1]
        slide = self.generator.presentation.slides.add_slide(slide_layout)
        
        from pptx.util import Inches
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        paragraph = textbox.text_frame.paragraphs[0]
        paragraph.text = "测试标题"
        
        # 测试不同级别的标题格式化
        for level in range(1, 7):
            self.generator._format_heading_by_level(paragraph, level)
            
            # 检查字体大小是否设置
            if paragraph.runs:
                font_size = paragraph.runs[0].font.size
                self.assertIsNotNone(font_size)
    
    def test_apply_design_suggestions(self):
        """测试设计建议应用"""
        self.generator._initialize_presentation()
        
        design_suggestions = {
            'color_scheme': {
                'primary': '#0078D4',
                'secondary': '#106EBE',
                'accent': '#005A9E'
            },
            'font_suggestions': {
                'title': 'Microsoft YaHei',
                'heading': 'Microsoft YaHei',
                'body': 'Microsoft YaHei'
            }
        }
        
        # 这个方法主要是框架，不会抛出异常即可
        try:
            self.generator._apply_design_suggestions(design_suggestions)
        except Exception as e:
            self.fail(f"应用设计建议时发生异常: {e}")
    
    def test_create_content_textbox(self):
        """测试内容文本框创建"""
        self.generator._initialize_presentation()
        
        # 创建测试幻灯片
        slide_layout = self.generator.presentation.slide_layouts[1]
        slide = self.generator.presentation.slides.add_slide(slide_layout)
        
        textbox = self.generator._create_content_textbox(slide)
        
        self.assertIsNotNone(textbox)
        self.assertTrue(hasattr(textbox, 'text_frame'))
    
    def test_format_table_cells(self):
        """测试表格单元格格式化"""
        self.generator._initialize_presentation()
        
        # 创建测试幻灯片和表格
        slide_layout = self.generator.presentation.slide_layouts[1]
        slide = self.generator.presentation.slides.add_slide(slide_layout)
        
        from pptx.util import Inches
        table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2))
        table = table_shape.table
        
        # 测试表头格式化
        header_cell = table.cell(0, 0)
        header_cell.text = "表头"
        self.generator._format_table_header_cell(header_cell)
        
        # 测试数据单元格格式化
        data_cell = table.cell(1, 0)
        data_cell.text = "数据"
        self.generator._format_table_data_cell(data_cell)
        
        # 检查格式是否应用（不会抛出异常即可）
        self.assertEqual(header_cell.text, "表头")
        self.assertEqual(data_cell.text, "数据")

if __name__ == '__main__':
    unittest.main()